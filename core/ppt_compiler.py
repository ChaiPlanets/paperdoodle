# core/ppt_compiler.py
import os
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Default Fallback Palette (Used only if template is missing)
COLOR_BG_DARK = RGBColor(15, 23, 42)      # Slate 900
COLOR_PRIMARY = RGBColor(30, 41, 59)      # Slate 800
COLOR_ACCENT = RGBColor(37, 99, 235)      # Blue 600
COLOR_MUTED = RGBColor(100, 116, 139)     # Slate 500
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT = RGBColor(51, 65, 85)         # Slate 700


def clean_markdown_text(text: str) -> str:
    """Strips markdown bold/italic tags and heading markers."""
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)
    return text.strip()


def extract_key_points(summary_text: str, max_points: int = 4) -> list[str]:
    """Extracts clean bullet points from markdown text."""
    lines = summary_text.splitlines()
    points = []
    
    for line in lines:
        cleaned = line.strip()
        if cleaned.startswith(("-", "*", "•", "1.", "2.", "3.", "4.")):
            clean_pt = re.sub(r"^[-*•\d.]+\s*", "", cleaned)
            clean_pt = clean_markdown_text(clean_pt)
            if len(clean_pt) > 15:
                points.append(clean_pt)
        elif cleaned and not cleaned.startswith("#") and len(cleaned) > 40:
            points.append(clean_markdown_text(cleaned))
            
    if not points:
        sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", summary_text) if len(s.strip()) > 20]
        points = sentences[:max_points]
        
    return points[:max_points]


def add_slide_header(slide, title_text: str, category: str = "RESEARCH BRIEFING"):
    """Adds a structured header to content slides."""
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(8.4), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_ACCENT

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.70), Inches(8.4), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY


def resolve_template_path() -> str | None:
    """Finds theme.pptx across standard locations with case-insensitivity."""
    candidates = [
        "Theme.pptx",
        "theme.pptx",
        os.path.join("assets", "Theme.pptx"),
        os.path.join("assets", "theme.pptx"),
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    return None


def create_deck(
    paper_title: str,
    summary_text: str,
    analogy_text: str,
    image_path: str,
    output_path: str
):
    """Compiles a 4-slide presentation adopting the template's master designs."""
    template_path = resolve_template_path()
    has_template = template_path is not None

    if has_template:
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    # Pick layout references
    title_layout = prs.slide_layouts[0]  # Standard Title Slide Layout
    content_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Preserves Template Background & Style)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(title_layout)

    if not has_template:
        # Only draw fallback dark box if NO template file exists
        bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = COLOR_BG_DARK
        bg1.line.color.rgb = COLOR_BG_DARK

    # Check if native placeholders exist on the template's title slide
    if len(slide1.placeholders) >= 2:
        title_placeholder = slide1.placeholders[0]
        subtitle_placeholder = slide1.placeholders[1]
        
        title_placeholder.text = clean_markdown_text(paper_title)
        subtitle_placeholder.text = f"Research Synthesis • {datetime.now().strftime('%B %d, %Y')}"
    elif len(slide1.placeholders) == 1:
        slide1.placeholders[0].text = clean_markdown_text(paper_title)
    else:
        # Fallback textbox overlay if layout has no placeholders
        title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8.0), Inches(2.5))
        tf1 = title_box.text_frame
        tf1.word_wrap = True

        p_tag = tf1.paragraphs[0]
        p_tag.text = "RESEARCH SYNTHESIS & BRIEFING"
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_ACCENT if has_template else RGBColor(96, 165, 250)
        p_tag.space_after = Pt(12)

        p_title = tf1.add_paragraph()
        p_title.text = clean_markdown_text(paper_title)
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY if has_template else COLOR_WHITE
        p_title.space_after = Pt(16)

        p_meta = tf1.add_paragraph()
        p_meta.text = f"Generated by PaperDoodle • {datetime.now().strftime('%B %d, %Y')}"
        p_meta.font.size = Pt(10)
        p_meta.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 2: Core Executive Summary
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(content_layout)
    add_slide_header(slide2, "Executive Overview & Problem Statement", "SUMMARY")

    points = extract_key_points(summary_text, max_points=4)
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.6))
    tf2 = content_box.text_frame
    tf2.word_wrap = True

    for i, pt in enumerate(points):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"•   {pt}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(14)
        p.line_spacing = 1.2

    # -------------------------------------------------------------
    # SLIDE 3: Visual & Mechanism (Side-by-Side)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(content_layout)
    add_slide_header(slide3, "Visual Architecture & Mechanism", "DIAGRAM & INTUITION")

    if image_path and os.path.exists(image_path):
        slide3.shapes.add_picture(
            image_path,
            Inches(0.8),
            Inches(1.4),
            width=Inches(4.1)
        )

    right_box = slide3.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.0), Inches(3.6))
    tf3 = right_box.text_frame
    tf3.word_wrap = True

    p_h = tf3.paragraphs[0]
    p_h.text = "How It Works"
    p_h.font.size = Pt(15)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_PRIMARY
    p_h.space_after = Pt(8)

    clean_analogy = clean_markdown_text(analogy_text)
    if len(clean_analogy) > 420:
        clean_analogy = clean_analogy[:417] + "..."

    p_body = tf3.add_paragraph()
    p_body.text = clean_analogy
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = COLOR_TEXT
    p_body.line_spacing = 1.2

    # -------------------------------------------------------------
    # SLIDE 4: Strategic Takeaways & Impact
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(content_layout)
    add_slide_header(slide4, "Strategic Takeaways & Significance", "OUTCOMES")

    all_points = extract_key_points(summary_text, max_points=8)
    takeaway_points = all_points[2:6] if len(all_points) > 4 else all_points[:3]

    content_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.6))
    tf4 = content_box4.text_frame
    tf4.word_wrap = True

    for i, pt in enumerate(takeaway_points):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.text = f"•   {pt}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(14)
        p.line_spacing = 1.2

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path